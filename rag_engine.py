import os
import uuid
import requests
from dotenv import load_dotenv
import chromadb
from chromadb.utils import embedding_functions
from google import genai
from pypdf import PdfReader

load_dotenv("C:\\nlp_6th_sem\\venv\\.env")

# Initialize vector DB locally
CHROMA_DATA_PATH = "./chroma_db"
COLLECTION_NAME = "astro_data"

# We use PersistentClient so the database saves to the hard drive
client = chromadb.PersistentClient(path=CHROMA_DATA_PATH)

# Using a robust but lightweight embedding function (all-MiniLM-L6-v2 runs fast locally on CPU)
sentence_transformer_ef = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")

# Get or create the collection for storing our Astro data
collection = client.get_or_create_collection(
    name=COLLECTION_NAME, 
    embedding_function=sentence_transformer_ef
)

def process_document(file_path: str, filename: str) -> int:
    """Read a document, extract text, split into chunks, and embed into ChromaDB."""
    text = ""
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        reader = PdfReader(file_path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted: text += extracted + "\n"
    elif ext == 'csv':
        import csv
        with open(file_path, newline='', encoding='utf-8', errors='ignore') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                text += ", ".join(row) + "\n"
    elif ext in ['pptx', 'ppt']:
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    else:
        raise ValueError(f"Unsupported file format: {ext}")
            
    # Basic text chunking strategy
    chunk_size = 1000
    overlap = 200
    chunks = []
    
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
        
    # Generate unique IDs and metadata for each chunk
    ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
    metadatas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
    
    # Batch add to ChromaDB
    if chunks:
        collection.add(
            documents=chunks,
            metadatas=metadatas,
            ids=ids
        )
        
    return len(chunks)

def query_astro_gpt(user_query: str, history: list = []):
    """Retrieve relevant chunks from DB and prompt the LLM API, with NASA fallback."""
    gemini_api_key = os.getenv("gemini_api")
    
    # 1. Retrieve the top 3 most relevant context chunks using our local database
    results = collection.query(
        query_texts=[user_query],
        n_results=3
    )
    
    context = ""
    retrieved_chunks = []
    
    # Simple check if meaningful documents were returned
    if results['documents'] and results['documents'][0]:
        retrieved_chunks = results['documents'][0]
        context = "\n\n---\n\n".join(retrieved_chunks)
        
    # 2. NASA API FALLBACK
    if not context or context.strip() == "":
        try:
            nasa_res = requests.get(f"https://images-api.nasa.gov/search?q={user_query}&media_type=image")
            if nasa_res.status_code == 200:
                nasa_data = nasa_res.json()
                items = nasa_data.get("collection", {}).get("items", [])
                if items:
                    first_item = items[0]["data"][0]
                    title = first_item.get("title", "No Title")
                    description = first_item.get("description", "No Description Available.")
                    context = f"[NASA API FALLBACK]\nTitle: {title}\nDescription: {description}"
                    retrieved_chunks = [f"[NASA Database]: {title}"]
                else:
                    context = "No relevant context found in local DB or NASA public archives."
        except Exception as e:
            context = "No local context found. NASA API fallback unavailable."
        
    # 3. Build the System Prompt tailored for Astro-GPT
    prompt = f"""You are Astro-GPT, an expert Small Language Model and brilliant interpreter for astronomical data, astrophysics, and space science.
You have been provided with the following retrieved contextual information from domain-specific data or NASA.
Use this context to accurately and intelligently answer the user's question. 

Context Database:
{context}

User Question: {user_query}
"""
    
    # 4. Format the Gemini History payload
    gemini_contents = []
    
    # Reconstruct the trailing conversation log to be consumed natively via the content queue
    for msg in history:
        if 'role' not in msg or 'content' not in msg: continue
        # Gemini expects standard alternating text dict streams formatted as below
        role = "model" if msg['role'] == "assistant" else "user"
        gemini_contents.append({"role": role, "parts": [{"text": msg['content']}]})
        
    # Append the contextualized final user query as the current frame
    gemini_contents.append({"role": "user", "parts": [{"text": prompt}]})
    
    # 5. Call the new Gemini GenAI client
    client = genai.Client(api_key=gemini_api_key)
    
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=gemini_contents
    )
    
    return response.text, retrieved_chunks
